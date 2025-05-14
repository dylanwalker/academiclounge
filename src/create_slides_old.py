import os
import pptx
import copy
#from pptx.oxml import parse_xml
#from pptx.oxml.ns import qn

# Path to presentations
path_to_presentations = os.path.join(os.path.dirname(__file__), os.pardir, "slides")
path_to_output = os.path.join(os.path.dirname(__file__), os.pardir, "slideshow")
if not os.path.exists(path_to_output):
    os.makedirs(path_to_output)


# Create a new presentation
merged_presentation = pptx.Presentation()

merged_presentation.slides

pptx_files = [os.path.join(path_to_presentations, f) for f in os.listdir(path_to_presentations) if f.endswith('.pptx')]

for pptx_file in pptx_files:
    curr_presentation = pptx.Presentation(pptx_file)
    for slide in curr_presentation.slides:
        slide_layout =  slide.slide_layout # merged_presentation.slide_layouts[5]#slide.slide_layout
        curr_slide = merged_presentation.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        print(slide_layout)
        #copied_slide = duplicate_slide(curr_presentation, slide_ind)
        #merged_presentation.slides.add_slide(copied_slide.slide_layout)
    print(f"Finished processing {pptx_file}")

merged_presentation.save(os.path.join(path_to_output, "merged_presentation.pptx"))

print(f"Saved merged presentation to {os.path.join(path_to_output, 'merged_presentation.pptx')}")
